"""
app.py — FastAPI application for the Bid Desk AI service.

Provides:
  - Lifespan startup gate: calls verify_access() on boot; aborts loudly with a
    clear error if the org/key lacks gpt-5.4 / gpt-5.4-mini access (D-16, LOCKED).
    Trade-off: a transient OpenAI outage prevents server start — correct behaviour
    for a graded prototype where access proof matters more than resilience. No bypass
    toggle (CLAUDE.md §2: no unrequested escape hatches).

  - GET /health: hermetic liveness probe (the docker-compose healthcheck target).
    Returns {"status": "ok"} with no model call.

  - GET /stream/demo: streams the trivial LangGraph demo graph as SSE events
    in {type, payload} format, observable via `curl -N` (PLAT-04).

  - GET /data/rfq: live-regenerate the marketing-services RFQ (DATA-04).

  - POST /data/vendor-gen: live-regenerate a vendor response for a given persona (DATA-04).
    Accepts optional rfq_text body parameter so callers can pass the same RFQ context
    to all vendors (ensuring a valid comparison). If rfq_text is omitted, a fresh RFQ
    is generated inline.

  - POST /extract/file-text: accept file upload (PDF/DOCX/XLSX/PPTX), extract text
    best-effort, return {text, filename, chars} (D-05, INPUT-01).

  - POST /input/raw-text: wrap raw text + vendor name into a VendorResponse JSON
    (D-06, INPUT-02).

  - CORS: localhost:3000 + all *.vercel.app preview/prod URLs (SHIP-01).

  - X-Accel-Buffering: no header unconditionally on SSE endpoints so streams
    pass through Render's nginx proxy without buffering (SHIP-01).

Security: verify_access() never logs the API key (T-03-01 mitigation).
POST /data/vendor-gen validates persona against MESS_SPECS keys before use (T-02-11).
POST /extract/file-text: filename is used only for extension detection; raw bytes
never interpreted as a path (T-05-02-A). Files > 20 MB are rejected with 413 (T-05-02-B).
POST /input/raw-text: raw_text capped at 200_000 chars via pydantic_Field (T-05-02-C).
CORS: no wildcard origin; allow_origin_regex used for Vercel subdomain matching (T-05-02-D).
"""

from __future__ import annotations

import io
from collections.abc import AsyncGenerator
from contextlib import asynccontextmanager

import docx
import openpyxl
import pptx
import pypdf
from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.concurrency import run_in_threadpool
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, model_validator
from pydantic import Field as pydantic_Field
from sse_starlette import EventSourceResponse

from agents._demo import demo_graph
from agents.comparison import comparison_graph
from agents.extraction import extraction_graph
from agents.rfq_gen import generate_rfq, render_rfq_md
from agents.vendor_gen import MESS_SPECS, generate_vendor_response
from llm.factory import verify_access
from prompts.registry import load as load_prompt
from schemas.domain import RFQ, ExtractionResult, VendorResponse
from schemas.events import EventEnvelope


@asynccontextmanager
async def lifespan(app_: FastAPI) -> AsyncGenerator[None, None]:
    """FastAPI lifespan handler — calls verify_access() on boot (D-16).

    If verify_access() raises (access denied or param error), the exception
    propagates and uvicorn aborts startup with the clear failure message.
    This is the second half of D-16; the first half is scripts/verify_access.py.
    """
    verify_access()  # raises RuntimeError on missing access; aborts startup
    yield


app = FastAPI(title="Bid Desk AI", lifespan=lifespan)

# ponytail: allow_origins=["https://*.vercel.app"] is dead config — Starlette exact-matches
# allow_origins by string equality. allow_origin_regex handles subdomain wildcards correctly.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"],
    allow_origin_regex=r"https://.*\.vercel\.app",
    allow_methods=["GET", "POST"],
    allow_headers=["Content-Type"],
)


class VendorGenRequest(BaseModel):
    """Request body for POST /data/vendor-gen (DATA-04)."""

    persona: str = pydantic_Field(max_length=64)
    rfq_text: str | None = pydantic_Field(default=None, max_length=200_000)


class RawTextInput(BaseModel):
    """Request body for POST /input/raw-text (D-06, INPUT-02)."""

    vendor_name: str = pydantic_Field(max_length=200)
    raw_text: str = pydantic_Field(max_length=200_000)  # T-05-02-C: DoS guard


def _extract_text(content: bytes, suffix: str) -> str:
    """Best-effort text extraction from file bytes by extension (D-05).

    Security (T-05-02-A): suffix is derived from the filename extension in the
    endpoint — the raw filename is never passed here; path traversal is impossible.
    All parser errors are caught and return "" so the caller always gets a string.
    """
    try:
        if suffix == "pdf":
            reader = pypdf.PdfReader(io.BytesIO(content))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        if suffix == "docx":
            doc = docx.Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        if suffix == "xlsx":
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            parts: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            parts.append(str(cell.value))
            return "\n".join(parts)
        if suffix == "pptx":
            prs = pptx.Presentation(io.BytesIO(content))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            parts.append(para.text)
            return "\n".join(parts)
    except Exception:  # noqa: BLE001
        pass
    return ""


@app.post("/extract/file-text")
async def extract_file_text(file: UploadFile = File(...)) -> dict:
    """Accept a file upload and return extracted plain text (D-05, INPUT-01).

    Supports PDF, DOCX, XLSX, PPTX via best-effort extraction (full OCR not required,
    assignment §11). Returns {text, filename, chars} regardless of extraction quality —
    low yield is a quality signal the UI surfaces, not a server error.

    Security: filename used only for extension detection (T-05-02-A).
    Files > 20 MB rejected with 413 before reading into memory (T-05-02-B).
    """
    content = await file.read()
    if len(content) > 20_000_000:  # T-05-02-B: 20 MB app-layer cap (413)
        raise HTTPException(status_code=413, detail="File too large (>20 MB)")
    suffix = (file.filename or "").rsplit(".", 1)[-1].lower()
    text = _extract_text(content, suffix)
    return {"text": text, "filename": file.filename, "chars": len(text)}


@app.post("/input/raw-text")
async def wrap_raw_text(req: RawTextInput) -> dict:
    """Wrap raw pasted text + vendor name into a VendorResponse JSON (D-06, INPUT-02).

    Output is a valid VendorResponse so it feeds directly into /extract/vendor
    without any schema drift. The extraction agent reads raw_text and produces
    ExtractionResult — no pre-extracted fields are added here.
    """
    vendor = VendorResponse(
        vendor_name=req.vendor_name,
        persona="buyer-upload",
        mess_spec=[],
        source_id=f"upload-{req.vendor_name[:20]}",
        format_label="text",
        raw_text=req.raw_text,
    )
    return vendor.model_dump(mode="json")


@app.get("/health")
async def health() -> dict:
    """Liveness probe — the docker-compose healthcheck target. Makes no model call.

    Hermetic: returns a static body so it can run before/without OpenAI access
    (the verify_access boot gate is a separate uvicorn-startup concern).
    """
    return {"status": "ok"}


@app.get("/prompts/{prompt_id}")
async def get_prompt(prompt_id: str) -> dict:
    """Return a Prompt Pack entry (metadata + markdown body) by id.

    The Prompt Trace screen fetches this on demand so the prompts have a single source
    of truth — the versioned files in services/ai/prompts/ — instead of being copied into
    the web app. registry.load() validates the id (^[a-z0-9-]+$, path-traversal safe) and
    resolves the latest version.

    # ponytail: file-backed via the registry today. Moving the Prompt Pack into a database
    # is a deferred option — keep this endpoint as the single read seam if/when that lands.
    """
    try:
        post = load_prompt(prompt_id)
    except ValueError as exc:  # invalid id shape
        raise HTTPException(status_code=400, detail=str(exc)) from exc
    except KeyError as exc:  # no such prompt
        raise HTTPException(status_code=404, detail=f"No prompt '{prompt_id}'") from exc
    return {"id": prompt_id, "metadata": post.metadata, "content": post.content}


@app.get("/data/rfq")
async def get_rfq() -> dict:
    """Live-regenerate the marketing-services RFQ via rfq-gen prompt (DATA-04).

    Makes a live OpenAI call — not served from the committed fixture.
    Returns: JSON-serializable RFQ dict.

    generate_rfq() is a blocking (synchronous) LangChain call; run it in the threadpool
    so it never blocks the event loop. Otherwise a single in-flight RFQ regen stalls the
    /health probe (→ container marked unhealthy) and every other concurrent request.
    """
    rfq = await run_in_threadpool(generate_rfq)
    return rfq.model_dump(mode="json")


@app.post("/data/vendor-gen")
async def post_vendor_gen(req: VendorGenRequest) -> dict:
    """Live-regenerate a vendor response for the given persona (DATA-04).

    req.persona: one of "thorough-but-pricey", "cheap-but-incomplete", "polished-fluff".
    req.rfq_text: optional RFQ Markdown text. If provided, all vendors see the same RFQ
                  (required for a valid comparison). If omitted, a fresh RFQ is generated.
    Returns: JSON-serializable VendorResponse dict.

    Security: persona is validated against MESS_SPECS keys before use (T-02-11).
    """
    if req.persona not in MESS_SPECS:
        raise HTTPException(
            status_code=400,
            detail=f"Unknown persona: {req.persona!r}. Must be one of {list(MESS_SPECS)}",
        )
    rfq_text = req.rfq_text
    if rfq_text is None:
        rfq = await run_in_threadpool(generate_rfq)
        rfq_text = render_rfq_md(rfq)
    # Blocking LangChain call — offload so it never blocks the event loop (see get_rfq).
    vendor = await run_in_threadpool(
        generate_vendor_response, rfq_text, req.persona, MESS_SPECS[req.persona]
    )
    return vendor.model_dump(mode="json")


class ExtractionRequest(BaseModel):
    """Request body for POST /extract/vendor (EXTRACT-03).

    vendor_response carries the full VendorResponse; rfq carries the full RFQ.

    # ponytail: max_length on nested str fields can't be enforced via
    # pydantic_Field (it only applies to direct str fields, not nested model
    # attributes), so a model_validator caps both the vendor raw_text and the
    # aggregate request size. The aggregate cap covers the RFQ free-text fields
    # (questionnaire, compliance_requirements, line-item descriptions/deliverables)
    # that an attacker could otherwise inflate independently of raw_text (W-R2).
    # This is an app-layer best-effort cap — the body is already deserialized by
    # the time it runs. The authoritative request-body size limit is server-level
    # (uvicorn --limit-* / reverse proxy) and is deferred to SHIP-01.
    """

    vendor_response: VendorResponse
    rfq: RFQ

    _MAX_RAW_TEXT = 200_000
    _MAX_TOTAL = 500_000

    @model_validator(mode="after")
    def _check_payload_size(self) -> ExtractionRequest:
        if len(self.vendor_response.raw_text) > self._MAX_RAW_TEXT:
            raise ValueError(
                f"vendor_response.raw_text exceeds {self._MAX_RAW_TEXT:,} chars "
                f"(got {len(self.vendor_response.raw_text)})"
            )
        total = len(self.model_dump_json())
        if total > self._MAX_TOTAL:
            raise ValueError(
                f"extraction request exceeds {self._MAX_TOTAL:,} chars total "
                f"(got {total}) — RFQ free-text included"
            )
        return self


@app.get("/stream/demo")
async def stream_demo() -> EventSourceResponse:
    """Stream the trivial LangGraph demo graph as SSE events (PLAT-04).

    Emits the full EVENT_TYPES taxonomy in order:
      data: {"type":"status","payload":{...}}
      data: {"type":"partial","payload":{...}}
      data: {"type":"result","payload":{...}}
      data: {"type":"done","payload":{}}

    Observable via: curl -N http://localhost:8000/stream/demo
    No model calls on this path — the demo graph emits hardcoded taxonomy events.
    """

    async def _generate() -> AsyncGenerator[dict, None]:
        # Validate every chunk through the closed envelope before serializing —
        # a malformed emit (unknown type, missing payload, extra keys) fails
        # loudly here rather than streaming unchecked to the client.
        async for chunk in demo_graph.astream({}, stream_mode="custom"):
            yield {"data": EventEnvelope(**chunk).model_dump_json()}
        # Final "done" event appended by the route after the graph completes.
        yield {"data": EventEnvelope(type="done", payload={}).model_dump_json()}

    return EventSourceResponse(_generate())


class ComparisonRequest(BaseModel):
    """Request body for POST /compare/vendors (COMPARE-01..05).

    Accepts a list of ExtractionResult objects + the original RFQ. The comparison
    agent never reads raw vendor text — only validated ExtractionResult[] (COMPARE-01).

    Vendor count is limited to _MAX_VENDORS (prototype limit — single-call context
    window constraint; RESEARCH Pitfall 7 / Review Fix LOW).

    Note: no char cap needed since ExtractionResult has no raw_text; vendor count guard only.

    Security: clamp runs server-side before result event (D-03 / Review Fix 9);
    exactly one result event (Review Fix 9).
    """

    extractions: list[ExtractionResult]
    rfq: RFQ

    _MIN_VENDORS: int = 2
    _MAX_VENDORS: int = 5
    # ponytail: _MAX_VENDORS=5 is a prototype limit — single-call context window constraint
    # (RESEARCH Pitfall 7 / Review Fix LOW). Increase if multi-vendor truncation is observed.

    @model_validator(mode="after")
    def _check_vendor_count(self) -> ComparisonRequest:
        n = len(self.extractions)
        if n < self._MIN_VENDORS:
            # A comparison needs at least two vendors. Reject the degenerate case
            # at the trust boundary instead of returning a semantically empty
            # ComparisonResult that silently hides the absence. (Review CR-02 / CLAUDE.md §1)
            raise ValueError(
                f"Too few vendors: {n} < {self._MIN_VENDORS}. "
                f"A comparison requires at least {self._MIN_VENDORS} vendors."
            )
        if n > self._MAX_VENDORS:
            raise ValueError(
                f"Too many vendors: {n} > {self._MAX_VENDORS} (prototype limit). "
                f"Submit at most {self._MAX_VENDORS} vendors per comparison request."
            )
        return self


@app.post("/compare/vendors")
async def compare_vendors(req: ComparisonRequest) -> EventSourceResponse:
    """Stream vendor comparison as SSE events (COMPARE-01..05).

    Accepts list[ExtractionResult] + RFQ. Runs the 4-node comparison graph:
    align → comparability → compare → clarify.

    Clamp runs server-side before the result event (D-03 / Review Fix 9).
    Exactly one result event is emitted — after clamp + clarification (Review Fix 9).
    All structured-output failure shapes emit safe error events.
    """

    async def _generate() -> AsyncGenerator[dict, None]:
        async for chunk in comparison_graph.astream(
            {"extractions": req.extractions, "rfq": req.rfq},
            stream_mode="custom",
        ):
            yield {"data": EventEnvelope(**chunk).model_dump_json()}
        # The clarify node owns the terminal `done` event (both its error and
        # success paths emit exactly one). Do NOT append another here. (Review CR-01)

    # ponytail: setting header in code is authoritative; env var may also help but is platform-specific
    response = EventSourceResponse(_generate())
    response.headers["X-Accel-Buffering"] = "no"
    return response


@app.post("/extract/vendor")
async def extract_vendor(req: ExtractionRequest) -> EventSourceResponse:
    """Stream vendor extraction as SSE events (EXTRACT-03).

    Accepts a VendorResponse + RFQ, runs the extraction graph, and streams
    status -> result -> done events. All structured-output failure shapes emit
    a safe error event (B-R2). Grounding runs before the result event (D-07).
    """

    async def _generate() -> AsyncGenerator[dict, None]:
        async for chunk in extraction_graph.astream(
            {"vendor": req.vendor_response, "rfq": req.rfq},
            stream_mode="custom",
        ):
            yield {"data": EventEnvelope(**chunk).model_dump_json()}
        yield {"data": EventEnvelope(type="done", payload={}).model_dump_json()}

    # ponytail: setting header in code is authoritative; env var may also help but is platform-specific
    response = EventSourceResponse(_generate())
    response.headers["X-Accel-Buffering"] = "no"
    return response
